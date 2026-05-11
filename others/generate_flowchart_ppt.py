"""生成定标流程与解调算法流程图 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# 颜色定义
TEAL = RGBColor(0x0F, 0x76, 0x6E)
TEAL_LIGHT = RGBColor(0xCC, 0xFB, 0xF1)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xF3, 0xF4, 0xF6)
BORDER = RGBColor(0xD1, 0xD5, 0xDB)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)


def add_box(slide, left, top, width, height, text, fill_color=WHITE,
            border_color=TEAL, font_size=10, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """添加一个带文字的方框"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = DARK
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    return shape


def add_diamond(slide, left, top, width, height, text, fill_color=TEAL_LIGHT):
    """添加菱形判断框"""
    return add_box(slide, left, top, width, height, text,
                   fill_color=fill_color, shape_type=MSO_SHAPE.DIAMOND, font_size=9)


def add_arrow(slide, start_x, start_y, end_x, end_y, label=""):
    """添加带箭头的连接线"""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = GRAY
    connector.line.width = Pt(1.2)
    # 箭头
    connector.begin_x = start_x
    connector.begin_y = start_y
    connector.end_x = end_x
    connector.end_y = end_y
    return connector


def add_text_label(slide, left, top, text, font_size=8, color=GRAY):
    """添加文字标签"""
    txBox = slide.shapes.add_textbox(left, top, Cm(2), Cm(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"
    return txBox


def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ========== 第1页：定标阶段流程图 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 标题
    title = slide1.shapes.add_textbox(Cm(1), Cm(0.3), Cm(20), Cm(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "定标阶段流程图"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = TEAL
    run.font.name = "Microsoft YaHei"

    # 流程框参数
    bw = Cm(4.5)   # box width
    bh = Cm(1.5)   # box height
    gap_v = Cm(0.8)  # vertical gap
    start_x = Cm(4)
    start_y = Cm(2)

    # 1. 系统初始化
    y = start_y
    add_box(slide1, start_x, y, bw, bh, "系统初始化\n连接激光器/VNA/Zynq",
            fill_color=TEAL_LIGHT, font_size=10, bold=True)

    # 2. 设置MRR电压
    y += bh + gap_v
    add_box(slide1, start_x, y, bw, bh, "设置 MRR 电压\nV = V_ref (1.25V)",
            font_size=10)

    # 3. 设置激光器波长
    y += bh + gap_v
    add_box(slide1, start_x, y, bw, bh, "设置激光器波长\nλ = λ_ref + Δλᵢ/1000",
            font_size=10)

    # 4. 等待稳定
    y += bh + gap_v
    add_box(slide1, start_x, y, bw, bh, "等待波长稳定\n(1 秒)", font_size=10)

    # 5. VNA测量
    y += bh + gap_v
    add_box(slide1, start_x, y, bw, bh, "VNA 测量 S21\n6001点, 10MHz~30GHz",
            font_size=10)

    # 6. 保存CSV
    y += bh + gap_v
    add_box(slide1, start_x, y, bw, bh, "保存 CSV 文件\ndelta_lambda_{Δλ}.csv",
            font_size=10)

    # 7. 判断
    y += bh + gap_v
    dw = Cm(4)
    dh = Cm(2)
    add_diamond(slide1, start_x + Cm(0.25), y, dw, dh, "所有 Δλ\n完成?")

    # 右侧分支：构建定标表
    rx = start_x + bw + Cm(3)
    ry = start_y + Cm(2)

    add_box(slide1, rx, ry, Cm(5), bh, "移动平均平滑 (N=5)\n提取谐振峰频率 f_peak",
            fill_color=SURFACE, font_size=10)

    add_box(slide1, rx, ry + bh + gap_v, Cm(5), Cm(2),
            "构建定标映射表\nT = {(Δλᵢ, f_peak_i, Curve_i)}\n按 f_peak 升序排列",
            fill_color=TEAL_LIGHT, font_size=10, bold=True)

    # 参数说明框
    param_x = rx + Cm(6)
    add_box(slide1, param_x, start_y, Cm(6.5), Cm(5),
            "定标参数\n━━━━━━━━━━━━━━\n"
            "λ_ref = 1551.85 nm\n"
            "V_ref = 1.25 V\n"
            "Δλ 范围: -160 ~ +160 pm\n"
            "Δλ 步长: 2 pm (161条)\n"
            "VNA: 10MHz~30GHz, 6001点\n"
            "功率: -10 dBm\n"
            "IF BW: 1000 Hz",
            fill_color=SURFACE, border_color=BORDER, font_size=9)

    # 添加"否"和"是"标签
    add_text_label(slide1, start_x - Cm(1.5), y + Cm(0.7), "否 ↑", font_size=9, color=GRAY)
    add_text_label(slide1, start_x + dw + Cm(0.5), y + Cm(0.7), "是 →", font_size=9, color=TEAL)

    # ========== 第2页：解调阶段流程图 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    title2 = slide2.shapes.add_textbox(Cm(1), Cm(0.3), Cm(20), Cm(1.2))
    tf2 = title2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "解调阶段流程图（归一化互相关匹配法）"
    run2.font.size = Pt(24)
    run2.font.bold = True
    run2.font.color.rgb = TEAL
    run2.font.name = "Microsoft YaHei"

    # 解调流程
    sx = Cm(2)
    sy = Cm(2)
    bw2 = Cm(5)
    bh2 = Cm(1.5)
    gap = Cm(0.7)

    # 步骤1
    y = sy
    add_box(slide2, sx, y, bw2, bh2,
            "① 加载实测 S21 曲线\n(多电压扫描数据)",
            fill_color=TEAL_LIGHT, font_size=10, bold=True)

    # 步骤2
    y += bh2 + gap
    add_box(slide2, sx, y, bw2, bh2,
            "② 提取实测峰值频率\nf_meas = argmax M̃(f)",
            font_size=10)

    # 步骤3
    y += bh2 + gap
    add_box(slide2, sx, y, bw2, bh2,
            "③ 二分查找 K 条候选\n定标曲线 (K=3)",
            font_size=10)

    # 步骤4
    y += bh2 + gap
    add_box(slide2, sx, y, bw2, bh2,
            "④ 频率轴插值对齐\nΔf = 0.01 GHz, 线性插值",
            font_size=10)

    # 步骤5
    y += bh2 + gap
    add_box(slide2, sx, y, bw2, bh2,
            "⑤ 归一化互相关\nρ = Σ(a-ā)(b-b̄)/(n·σₐ·σᵦ)",
            font_size=10)

    # 步骤6
    y += bh2 + gap
    add_box(slide2, sx, y, bw2, bh2,
            "⑥ 选择全局 ρ_max\n对应 Δλ* 和 V*",
            fill_color=SURFACE, font_size=10)

    # 步骤7
    y += bh2 + gap
    add_box(slide2, sx, y, bw2, bh2,
            "⑦ 温度计算\nT = T₀ + Δλ*/α",
            fill_color=TEAL_LIGHT, font_size=10, bold=True)

    # 右侧：公式说明
    fx = sx + bw2 + Cm(2)
    add_box(slide2, fx, sy, Cm(7), Cm(3.5),
            "归一化互相关系数 (NCC)\n━━━━━━━━━━━━━━━━━━━━\n\n"
            "        Σᵢ (aᵢ - ā)(bᵢ - b̄)\n"
            "ρ = ─────────────────────\n"
            "          n · σₐ · σᵦ\n\n"
            "ρ ∈ [-1, 1]\n"
            "ρ → 1: 曲线形状高度相似\n"
            "ρ > 0.9: 高置信度匹配",
            fill_color=WHITE, border_color=TEAL, font_size=10)

    # 右侧：温度公式
    add_box(slide2, fx, sy + Cm(4.2), Cm(7), Cm(2.5),
            "温度解调公式\n━━━━━━━━━━━━━━━━━━━━\n\n"
            "Δλ = α · ΔT    (α = 9.08 pm/°C)\n"
            "T = T₀ + Δλ*/α  (T₀ = 20°C)\n"
            "λ_FBG = λ_ref + Δλ/1000 (nm)",
            fill_color=WHITE, border_color=ORANGE, font_size=10)

    # 右下：算法特点
    add_box(slide2, fx, sy + Cm(7.2), Cm(7), Cm(2.8),
            "算法特点\n━━━━━━━━━━━━━━━━━━━━\n"
            "✓ 全频段形状匹配，鲁棒性强\n"
            "✓ 归一化处理，对幅度变化不敏感\n"
            "✓ 二分查找预筛选，计算高效\n"
            "✓ 多电压扫描，提高可靠性",
            fill_color=SURFACE, border_color=BORDER, font_size=9)

    # ========== 第3页：整体系统框图 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    title3 = slide3.shapes.add_textbox(Cm(1), Cm(0.3), Cm(25), Cm(1.2))
    tf3 = title3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "FBG+MRR 温度传感系统 — 数据流框图"
    run3.font.size = Pt(24)
    run3.font.bold = True
    run3.font.color.rgb = TEAL
    run3.font.name = "Microsoft YaHei"

    # 定标数据流
    cy = Cm(2.5)
    add_text_label(slide3, Cm(1), cy - Cm(0.5), "【定标阶段】", font_size=11, color=TEAL)

    boxes_cal = [
        ("可调谐激光器\nλ = λ_ref + Δλ", Cm(1)),
        ("FBG\n(波长滤波)", Cm(6)),
        ("MRR\n(V = V_ref)", Cm(11)),
        ("VNA\n(S21 测量)", Cm(16)),
        ("定标映射表\nΔλ ↔ f_peak", Cm(21.5)),
    ]
    for text, x in boxes_cal:
        add_box(slide3, x, cy, Cm(4), Cm(1.8), text, font_size=9)

    # 箭头标签
    for i in range(len(boxes_cal) - 1):
        x1 = boxes_cal[i][1] + Cm(4)
        x2 = boxes_cal[i+1][1]
        mid_x = (x1 + x2) / 2 - Cm(0.5)
        add_text_label(slide3, Emu(int(mid_x)), cy + Cm(0.6), "→", font_size=14, color=TEAL)

    # 解调数据流
    dy = Cm(5.5)
    add_text_label(slide3, Cm(1), dy - Cm(0.5), "【解调阶段】", font_size=11, color=ORANGE)

    boxes_dem = [
        ("Zynq 电压控制\n(V 扫描)", Cm(1)),
        ("MRR\n(电压调谐)", Cm(6)),
        ("VNA\n(S21 测量)", Cm(11)),
        ("归一化互相关\n曲线匹配", Cm(16)),
        ("温度输出\nT = T₀+Δλ/α", Cm(21.5)),
    ]
    for text, x in boxes_dem:
        color = TEAL_LIGHT if "温度" in text else WHITE
        add_box(slide3, x, dy, Cm(4), Cm(1.8), text, fill_color=color, font_size=9)

    for i in range(len(boxes_dem) - 1):
        x1 = boxes_dem[i][1] + Cm(4)
        x2 = boxes_dem[i+1][1]
        mid_x = (x1 + x2) / 2 - Cm(0.5)
        add_text_label(slide3, Emu(int(mid_x)), dy + Cm(0.6), "→", font_size=14, color=ORANGE)

    # 定标表到匹配的虚线连接
    add_text_label(slide3, Cm(18.5), Cm(4.3), "↓ 定标表", font_size=9, color=GRAY)

    # 保存
    output_path = "others/定标流程与解调算法_流程图.pptx"
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")


if __name__ == "__main__":
    create_ppt()

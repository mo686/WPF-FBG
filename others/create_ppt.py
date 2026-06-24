"""
开题报告 PPT 生成脚本
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BLACK = RGBColor(0x00, 0x00, 0x00)


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BLUE)
    # Title
    txBox = slide.shapes.add_textbox(Cm(2), Cm(6), Cm(22), Cm(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xBB, 0xDD, 0xFF)
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ACCENT_BLUE)
    txBox = slide.shapes.add_textbox(Cm(2), Cm(7), Cm(22), Cm(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Cm(1), Cm(0.3), Cm(23), Cm(1.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # Content
    txBox2 = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(22), Cm(14))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
        if bullet.startswith("  "):
            p.level = 1
            p.font.size = Pt(16)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    # ===== Slide 1: Title =====
    add_title_slide(
        prs,
        "应用于FBG传感的集成化解调方法\n及其温度补偿研究",
        "东南大学 · 专业学位硕士开题报告\n研究方向：光电集成与光传感"
    )

    # ===== Slide 2: Outline =====
    add_content_slide(prs, "汇报提纲", [
        "一、立题依据及价值",
        "二、研究内容及方法",
        "  2.1 基于参考FBG温度补偿的AWG解调系统",
        "  2.2 基于微波光子滤波器的FBG解调系统",
        "三、可行性分析与实验进展",
        "四、预期成果与工作计划",
    ])

    # ===== Section 1 =====
    add_section_slide(prs, "一、立题依据及价值")

    add_content_slide(prs, "研究背景", [
        "FBG传感器：抗电磁干扰、高灵敏度、体积小、易复用",
        "核心问题：精确解调反射中心波长的微小漂移",
        "AWG边缘滤波解调：结构简单、无机械扫描、多通道并行",
        "  2022年 北信科：1.6 nm动态范围，10 pm精度 (SOI)",
        "  2023年 之江实验室：C+L波段75 nm，优于2 pm",
        "  2024年 之江实验室：亚皮米级，RMSE 0.73 pm",
        "  课题组：双输入AWG + LASSO，34.45 nm范围，0.31 pm",
    ])

    add_content_slide(prs, "关键问题与研究思路", [
        "问题一：AWG温度漂移 → 解调误差",
        "  传统TEC方案功耗大、体积受限",
        "  LUT+插值法精度受限于标定密度（21.5 pm）",
        "  → 提出参考FBG实时温度补偿方案",
        "",
        "问题二：光域解调精度受限于光学滤波器分辨率",
        "  微波光子学：波长变化 → 微波频率变化",
        "  VNA频率分辨率可达Hz量级",
        "  → 建立微波光子滤波FBG解调系统",
    ])

    # ===== Section 2 =====
    add_section_slide(prs, "二、研究内容及方法")

    add_content_slide(prs, "2.1 参考FBG温度补偿的AWG解调", [
        "核心思想：",
        "  参考FBG（恒温）经AWG解调的表观波长变化",
        "  = AWG中心波长漂移量",
        "",
        "补偿算法：",
        "  Δλ_AWG = λ_ref_measured − λ_ref_true",
        "  λ_corrected = λ_measured − Δλ_AWG",
        "",
        "优势：仅需一路恒温参考通道，无需AWG精密温控",
    ])

    add_content_slide(prs, "2.1 系统链路设计", [
        "ASE宽带光源 → 光环形器 → 1×2光开关",
        "  FBG_ref（参考）/ FBG_sen（传感）分时切换",
        "反射光 → 光耦合器 → OSA监测 + AWG解调",
        "AWG → PD阵列 → 信号处理 → 差分补偿",
        "",
        "关键步骤：",
        "  ① FBG选型与安装（同封装、同温度环境）",
        "  ② 初始波长标定（OSA精确测量）",
        "  ③ 实时差分补偿解调",
    ])

    add_content_slide(prs, "2.2 微波光子滤波FBG解调系统", [
        "原理：PM-IM转换 + 微环drop端滤波",
        "  激光器 → PM → 微环drop端 → FBG反射",
        "  → EDFA → PD拍频 → VNA测S21",
        "",
        "微环谐振峰对准FBG中心波长时：",
        "  drop端损耗最小，VNA频域形成微波通带",
        "  FBG波长漂移 → 通带频率响应变化",
        "",
        "通过定标与匹配算法反推波长偏移量",
    ])

    add_content_slide(prs, "2.2 定标与解调算法", [
        "定标阶段：",
        "  固定微环电压，步进激光器波长（±160 pm, 1 pm步长）",
        "  采集321条S21曲线，提取峰值频率",
        "  构建 Δλ ↔ S21曲线 映射表",
        "",
        "解调算法（三级判决）：",
        "  ① 峰值频率预筛选（二分查找，排除90%+非匹配项）",
        "  ② 归一化互相关精细匹配（自适应频率窗口）",
        "  ③ 多功率全局择优（遍历所有扫描功率点）",
    ])

    add_content_slide(prs, "2.2 自适应频率窗口选取", [
        "目的：聚焦谐振峰主瓣，抑制带外噪声",
        "",
        "算法：",
        "  ① 提取实测曲线峰值频率 f_peak",
        "  ② 向两侧搜索3dB下降点，估算半功率半宽 Δf_3dB",
        "  ③ 匹配窗口 BW = 3 × Δf_3dB",
        "",
        "物理依据：",
        "  洛伦兹线型在 ±3γ 范围内包含 ~90% 峰值能量",
        "  超出此范围噪声主导，扩大窗口降低信噪比",
    ])

    # ===== Section 3 =====
    add_section_slide(prs, "三、可行性分析与实验进展")

    add_content_slide(prs, "3.1 AWG解调系统前期基础", [
        "芯片设计：",
        "  宽输入波导AWG（34 μm），3-dB带宽1.60 nm",
        "  双输入架构，有效通道间隔0.8 nm，范围34.45 nm",
        "",
        "算法：",
        "  多项式回归 + 高斯噪声增强 + 十折交叉验证",
        "  LASSO正则化，RMSE 0.31 pm，分辨率0.13 pm",
        "",
        "问题：未加温控时误差呈系统性负偏移（-6~0 pm）",
    ])

    add_content_slide(prs, "3.2 微波光子系统实验验证", [
        "FBG温度灵敏度标定：",
        "  25°C–50°C，α = 10.62 pm/°C，R² = 0.999",
        "",
        "温度传感验证（20°C → 25°C）：",
        "  最佳匹配 ρ = 0.9765，Δλ = 52 pm",
        "  温度误差 0.104°C",
        "",
        "多温度性能（0.1 mW步进）：",
        "  30°C–45°C范围内分辨率 < 1 pm",
        "  最优 0.2 pm @ 45°C（温度误差 0.02°C）",
    ])

    add_content_slide(prs, "3.2 自动化实验平台", [
        "PC上位机协调三台仪器：",
        "  VNA（GPIB）+ TSL激光器（GPIB/USB）+ Zynq FPGA（串口）",
        "",
        "软件架构（Python + PySide6）：",
        "  仪器驱动层 / 业务逻辑层 / GUI层",
        "",
        "可靠性设计：",
        "  电压自动置零保护",
        "  断电不丢失已采集数据",
        "  定标扫描支持暂停/恢复/中止",
    ])

    # ===== Section 4 =====
    add_section_slide(prs, "四、预期成果与工作计划")

    add_content_slide(prs, "预期成果", [
        "1. 基于参考FBG温度补偿的AWG解调系统",
        "  → 实现变温环境下的实时温度补偿",
        "",
        "2. 基于微环谐振器的微波光子FBG解调系统",
        "  → 定标-测量-互相关匹配完整方案",
        "  → 典型工况下分辨率 ≤ 1 pm",
        "",
        "3. 完成相关学术论文",
    ])

    add_content_slide(prs, "工作计划", [
        "2026.05–06  开题报告 + FBG应变测量 + PCB优化",
        "2026.07–08  AWG温度补偿实验系统搭建",
        "2026.09–10  AWG补偿多温度对比实验",
        "2026.11–12  提炼创新点，撰写学术论文",
        "2027.01–03  修改投稿，补充实验数据",
        "2027.04–05  完成硕士学位论文撰写",
        "2027.06     毕业答辩",
    ])

    # ===== End slide =====
    add_title_slide(prs, "谢谢！\n敬请指导", "")

    # Save
    output_path = r"others\开题报告PPT.pptx"
    prs.save(output_path)
    print(f"PPT saved: {output_path}")


if __name__ == "__main__":
    main()
